from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE

# Output path
OUT_PATH = "prezentace/cervena_karkulka.pptx"

prs = Presentation()
prs.slide_height = Inches(7.5)
prs.slide_width = Inches(13.333)

# Styles
TITLE_FONT = 'Poppins'
TITLE_SIZE = Pt(44)
BODY_FONT = 'Poppins'
BODY_SIZE = Pt(28)
PASTEL_RED = RGBColor(0xF6, 0xC7, 0xC7)
PASTEL_GREEN = RGBColor(0xCFF0D6)
PASTEL_BLUE = RGBColor(0xCFE8F5)
BG_COLOR = RGBColor(0xFF, 0xFB, 0xF0)

# Helper to add title + bullets
def add_slide(title, bullets, notes_text):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background - light cream
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Left illustration area (simple vector shapes)
    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.0), Inches(5.5), Inches(5.0))
    left.fill.solid()
    left.fill.fore_color.rgb = PASTEL_RED
    left.line.color.rgb = RGBColor(0xFF, 0xB3, 0xB3)

    # Title
    title_box = slide.shapes.add_textbox(Inches(6.3), Inches(0.6), Inches(6.4), Inches(1.2))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.name = TITLE_FONT
    title_p.font.size = TITLE_SIZE
    title_p.font.bold = True
    title_p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # Bullets box
    body_box = slide.shapes.add_textbox(Inches(6.3), Inches(1.6), Inches(6.4), Inches(4.6))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = body_tf.add_paragraph() if i>0 else body_tf.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = BODY_SIZE
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

    return slide

# Slides content (Czech)
slides = [
    (
        "Kdo je Červená Karkulka?",
        [
            "Malá dívka s červenou čepičkou a pláštěnkou.",
            "Bydlí s maminkou u lesa.",
            "Dnes nese dárek babičce."
        ],
        "Ukazujte velký obrázek Karkulky. Děti se mají seznámit s postavou. Ptejte se: Kdo nosí čepičku?"
    ),
    (
        "Dárek pro babičku",
        [
            "Maminka dala košík s koláčem a nápojem.",
            "Požádala: 'Jdi rovnou cestou a nezastavuj se.'",
            "Karkulka slíbila, že bude poslouchat."
        ],
        "Vysvětlete, proč je důležité poslouchat rodiče. Ukázat košík jako symbol péče."
    ),
    (
        "Setkání s vlkem",
        [
            "V lese potkala vlka.",
            "Vlk se ptá, kam jde, a tváří se přátelsky.",
            "Vlk má zlý plán — chce být mazaný."
        ],
        "Ukázat, že někdy se cizí lidé tváří přátelsky, ale nemusí to být dobře. Udržet tón přátelský, ne strašidelný."
    ),
    (
        "U babičky v domku",
        [
            "Vlk přišel k babičce dřív a schoval se v posteli.",
            "Karkulka si všimne, že 'babička' vypadá divně.",
            "Otázky: 'Proč máš tak velké oči/ucho/úst?'"
        ],
        "Použít bublinové otázky (dialog). Děti mohou odpovídat."
    ),
    (
        "Šťastný konec a co si zapamatovat",
        [
            "Myslivec přišel, pomohl a vlka porazil.",
            "Babička je v pořádku a všichni se radují.",
            "Poučení: poslouchej rodiče, choď po cestě a nemluv s cizími."
        ],
        "Závěrečné shrnutí a krátký úkol: Co bys udělal/la, kdyby tě zastavil cizinec?"
    )
]

for title, bullets, notes in slides:
    add_slide(title, bullets, notes)

prs.save(OUT_PATH)
print(f"Uloženo: {OUT_PATH}")
